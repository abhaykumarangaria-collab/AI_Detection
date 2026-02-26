from flask import Flask, request, render_template
from pptx import Presentation
import joblib
import re
import nltk
from nltk.corpus import stopwords
from werkzeug.utils import secure_filename
import os

# -------------------------
# INITIAL SETUP
# -------------------------

app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# Download once (will skip if already downloaded)
nltk.download("stopwords")

# -------------------------
# LOAD MODEL + VECTORIZER
# -------------------------
model = joblib.load("ppt_ai_model.pkl")
vectorizer = joblib.load("vectorizer.pkl")

stop_words = set(stopwords.words("english"))

# -------------------------
# PREPROCESS FUNCTION
# -------------------------
def preprocess(text):
    text = text.lower()
    text = re.sub(r"[^a-zA-Z\s]", "", text)
    words = text.split()
    words = [w for w in words if w not in stop_words]
    return " ".join(words)

# -------------------------
# EXTRACT TEXT FROM PPT
# -------------------------
def extract_text_from_ppt(file_path):
    presentation = Presentation(file_path)
    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text += shape.text + " "

    return text.strip()

# -------------------------
# HOME ROUTE
# -------------------------
@app.route("/")
def home():
    return render_template("index.html")

# -------------------------
# PREDICTION ROUTE
# -------------------------
@app.route("/predict_ppt", methods=["POST"])
def predict_ppt():

    if "file" not in request.files:
        return "No file uploaded!"

    file = request.files["file"]

    if file.filename == "":
        return "No selected file!"

    if not file.filename.endswith(".pptx"):
        return "Only .pptx files are allowed!"

    # Save file securely
    filename = secure_filename(file.filename)
    file_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
    file.save(file_path)

    try:
        # Extract text
        text = extract_text_from_ppt(file_path)

        if text.strip() == "":
            return "PPT contains no readable text."

        # Preprocess
        processed_text = preprocess(text)

        # Vectorize
        vector_text = vectorizer.transform([processed_text])

        # Predict
        prediction = model.predict(vector_text)[0]
        probability = model.predict_proba(vector_text)[0][1]

        result = "AI Generated PPT" if prediction == 1 else "Human Written PPT"
        probability = round(float(probability) * 100, 2)

    except Exception as e:
        return f"Error processing file: {str(e)}"

    finally:
        # Always remove file
        if os.path.exists(file_path):
            os.remove(file_path)

    return render_template("result.html", result=result, probability=probability)

# -------------------------
if __name__ == "__main__":
    app.run(debug=True)